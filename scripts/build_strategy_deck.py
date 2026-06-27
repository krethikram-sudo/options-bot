#!/usr/bin/env python3
"""Build the Outlay board/investor STRATEGY deck as a native .pptx.
Narrative form, derived from docs/market-analysis-ai-spend-governance.md and
docs/product-strategy.md. Speaker notes carry the talking track.
    python scripts/build_strategy_deck.py  ->  docs/outlay-strategy-deck.pptx
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PAPER=RGBColor(0xFA,0xF8,0xF3); INK=RGBColor(0x1F,0x24,0x30); MUT=RGBColor(0x5C,0x64,0x70)
GRN=RGBColor(0x0F,0x6B,0x4F); GRND=RGBColor(0x0B,0x51,0x3C); GRNL=RGBColor(0xE7,0xF1,0xEC)
AMBER=RGBColor(0x9A,0x5A,0x18); WHITE=RGBColor(0xFF,0xFF,0xFF); MINT=RGBColor(0xDD,0xF0,0xE7)
KICKM=RGBColor(0xCF,0xEE,0xDE); LINE=RGBColor(0xE1,0xDD,0xD2)
SERIF="Georgia"; SANS="Calibri"; LM=0.6; CW=12.13
ASSETS=os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),"docs","assets")
EMUW,EMUH=Inches(13.333),Inches(7.5)
prs=Presentation(); prs.slide_width,prs.slide_height=EMUW,EMUH; BLANK=prs.slide_layouts[6]

def slide(bg=PAPER):
    s=prs.slides.add_slide(BLANK); r=s.shapes.add_shape(1,0,0,EMUW,EMUH)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element); return s
def box(s,x,y,w,h):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tb.text_frame.word_wrap=True; return tb.text_frame
def para(tf,t,sz,col,*,bold=False,font=SANS,first=False,sa=6,align=PP_ALIGN.LEFT,bullet=False):
    p=tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(sa)
    if bullet: t="•  "+t
    r=p.add_run(); r.text=t; f=r.font; f.size=Pt(sz); f.bold=bold; f.name=font; f.color.rgb=col; return p
def notes(s,t): s.notes_slide.notes_text_frame.text=t.strip()
def kicker(s,t,col=GRN,y=0.5): para(box(s,LM,y,CW,0.5),t.upper(),13,col,bold=True,first=True)
def pill(s,t,x=LM,y=0.5,w=4.7,fill=GRNL,col=GRND):
    sh=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(0.55)); sh.fill.solid()
    sh.fill.fore_color.rgb=fill; sh.line.fill.background(); sh.shadow.inherit=False
    sh.text_frame.margin_top=Inches(0.05); para(sh.text_frame,t.upper(),12.5,col,bold=True,first=True)
def band(s,t,y,h=1.0,sz=23,fill=GRNL,col=GRND,x=LM,w=CW):
    sh=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h)); sh.fill.solid()
    sh.fill.fore_color.rgb=fill; sh.line.fill.background(); sh.shadow.inherit=False
    sh.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    para(sh.text_frame,t,sz,col,bold=True,font=SERIF,first=True,align=PP_ALIGN.CENTER)
def bullets(s,x,y,w,h,items,sz=16,gap=15,col=INK):
    tf=box(s,x,y,w,h)
    for i,t in enumerate(items): para(tf,t,sz,col,bullet=True,first=(i==0),sa=gap)
def card(s,x,y,w,h,title,body,*,border=LINE,fill=WHITE,bsz=14,gap=7,tcol=MUT):
    sh=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h)); sh.fill.solid()
    sh.fill.fore_color.rgb=fill; sh.line.color.rgb=border; sh.line.width=Pt(1.25); sh.shadow.inherit=False
    tf=sh.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.26); tf.margin_right=Inches(0.26)
    tf.margin_top=Inches(0.2); tf.vertical_anchor=MSO_ANCHOR.TOP
    para(tf,title.upper(),12.5,tcol,bold=True,first=True,sa=9)
    for ln in body: para(tf,ln[0],bsz,INK,bullet=True,sa=gap)
    return sh
def pic(s,name,x,y,w,border=True,bcolor=None):
    p=os.path.join(ASSETS,name); iw,ih=Image.open(p).size
    sh=s.shapes.add_picture(p,Inches(x),Inches(y),width=Inches(w))
    if border: sh.line.color.rgb=bcolor or LINE; sh.line.width=Pt(1.25)
    return w*ih/iw
def title(s,t,y=1.1,sz=32): para(box(s,LM,y,CW,0.9),t,sz,INK,bold=True,font=SERIF,first=True)

# 1 — cover
s=slide(GRN)
tf=box(s,LM,1.5,CW,1.7); p=tf.paragraphs[0]
for t,c in [("Outlay",WHITE),(".ai",KICKM)]:
    r=p.add_run(); r.text=t; r.font.size=Pt(64); r.font.bold=True; r.font.name=SERIF; r.font.color.rgb=c
para(box(s,LM,3.0,CW,0.7),"Put AI compute on a budget.",24,MINT,font=SERIF,first=True)
ln=s.shapes.add_shape(1,Inches(LM),Inches(4.0),Inches(3.4),Inches(0.045)); ln.fill.solid()
ln.fill.fore_color.rgb=KICKM; ln.line.fill.background(); ln.shadow.inherit=False
para(box(s,LM,4.35,CW,0.7),"Market & product strategy",22,WHITE,bold=True,first=True)
tf=box(s,LM,6.3,CW,1.0); para(tf,"Krethikram Gowrisankar · Founder",16,MINT,bold=True,first=True,sa=3)
para(tf,"outlay-ai.com",13,RGBColor(0xBF,0xE0,0xD0))
notes(s,"One-line frame: AI spend is the fastest-growing line item nobody can explain. We make it attributed, forecastable, and governable.")

# 2 — the shift
s=slide(); kicker(s,"The shift")
title(s,"AI spend is exploding — and moving to consumption.")
band(s,"88% of US enterprises are raising AI budgets · 62% by ≥10%",2.2,0.95,22)
card(s,LM,3.45,3.83,3.2,"Consumption, not seats",[("AI scales by tokens/compute, not users — per-seat can't price it.",)],bsz=15,border=GRN)
card(s,4.65,3.45,3.83,3.2,"Inference > training",[("Production usage overtook training — recurring cost, not one-time.",)],bsz=15,border=GRN)
card(s,8.8,3.45,3.93,3.2,"Code-gen is the breakout",[("~$4B, ~28% of enterprise LLM use; Claude ~42% dev share.",)],bsz=15,border=GRN)
notes(s,"Sources: Wharton/GBK 2025; Menlo Ventures 2025. The bill is big, growing, and now consumption-shaped — which is exactly what breaks visibility and seats.")

# 3 — the problem
s=slide(); kicker(s,"The problem")
title(s,"Consumption AI spend is opaque, multi-vendor, and unattributed.")
bullets(s,LM,2.4,4.5,4.2,[
 "Provider bills are aggregated at the account level — no team, ticket, or cost-center metadata.",
 "Spend is spread across OpenAI, Anthropic, Bedrock, Azure, Cursor, Copilot — no unified view.",
 "Finance can't map a credit to a workload; 78% of IT leaders saw surprise AI charges.",
 "Nobody can answer the question that matters: “What is AI costing us per team / feature / output?”",
],sz=15,gap=15)
pic(s,"diagram-architecture.png",5.35,2.5,7.6)
notes(s,"This is the gap. The question finance actually asks has no answer today. Sources: Flexera, Datadog, Vantage.")

# 4 — category inflection
s=slide(GRN); kicker(s,"The category is here",col=KICKM)
para(box(s,LM,1.7,CW,1.0),"From niche to near-universal in 24 months.",34,WHITE,bold=True,font=SERIF,first=True)
band(s,"98% of organizations now manage AI spend — up from 31% in 2024",3.0,1.1,24,fill=RGBColor(0x0B,0x51,0x3C),col=MINT)
tf=box(s,LM,4.5,CW,2.0)
para(tf,"The fastest cost-discipline adoption the FinOps survey has ever recorded.",18,MINT,first=True,sa=8)
para(tf,"Managing AI spend is the #1 FinOps priority for the next 12 months, and AI cost management is the #1 skill gap (58%).",18,MINT)
notes(s,"FinOps Foundation State of FinOps 2026 (n=1,192). There is now a budget line and an owner to sell to — the market arrived.")

# 5 — what Outlay is
s=slide(); kicker(s,"What Outlay is")
title(s,"The work-attribution + forecast + governance layer — not a dashboard.")
card(s,LM,2.35,3.83,4.3,"1 · Attribute",[("AI spend → the ticket, feature, engineer, and program — via the work-tracker join.",),("The question finance actually asks; hard to copy.",)],border=GRN,bsz=14)
card(s,4.65,2.35,3.83,4.3,"2 · Forecast",[("Back-tested on the customer's OWN delivered work (leave-one-out).",),("Their data, measured error — the trust unlock.",)],border=GRN,bsz=14)
card(s,8.8,2.35,3.93,4.3,"3 · Govern",[("Program budgets, real-time pacing, on-/off-track rating.",),("Turns a report into a system of record you act on.",)],border=GRN,bsz=14)
notes(s,"All three are built. The strategic move is positioning: sell the join + honest forecast + governance, not commodity 'visibility'. Read-only, BYOK, metadata-only.")

# 6 — defensibility
s=slide(); kicker(s,"Why it's defensible")
title(s,"Three layers compete; none owns ours.")
card(s,LM,2.35,3.83,4.3,"Visibility dashboards",[("Vantage · CloudZero · Finout · Datadog",),("Attribute at team/project — not the work-item join; crowded.",)],border=AMBER,bsz=13)
card(s,4.65,2.35,3.83,4.3,"LLM gateways",[("LiteLLM · Portkey · OpenRouter · TrueFoundry",),("In-path cost cuts + budgets — opposite of our read-only posture.",)],border=AMBER,bsz=13)
card(s,8.8,2.35,3.93,4.3,"Cloud optimizers",[("ProsperOps (Flexera) · Zesty · nOps",),("Reserve cloud GPUs you run — not the model-API layer; no work data.",)],border=AMBER,bsz=13)
band(s,"Our moat: work-item attribution · own-data forecast · metadata-only / read-only · budget governance",6.75,0.6,14)
notes(s,"We don't fight on 'visibility' or token price. The tracker join + honest forecast + read-only posture are what they don't lead with.")

# 7 — product proof
s=slide(); kicker(s,"The product today")
para(box(s,LM,1.05,CW,0.7),"Built and demoing — governance on real-shaped data.",26,INK,bold=True,font=SERIF,first=True)
pic(s,"shot-governance.png",2.57,1.95,8.2)
notes(s,"Proof it exists and works: program budgets, breach projections, earned-value ratings, the quarterly variance table. Live on sample data; verified with a real gov pilot motion (Maryland).")

# 8 — the wedge / ICP
s=slide(); kicker(s,"The wedge")
title(s,"Land the hardest, highest-value layer first.")
card(s,LM,2.35,6.0,4.3,"Beachhead ICP",[("Eng-heavy software / AI-native companies",),("250–5,000 employees",),(">$250–500k/yr AI spend, consumption-based, multi-provider, unowned",),("Buyer: platform/eng lead + VP Finance/FinOps",)],border=GRN,bsz=14,gap=10)
card(s,6.83,2.35,5.9,4.3,"Why here",[("The attribution gap is worst — and our tracker join only fires here",),("Budgets growing 50–100%/yr → urgency + a buyer with money",),("Reachable founder-led; no 12-month procurement",),("Gov/Maryland = parallel high-ACV lighthouse + compliance moat, NOT the roadmap driver",)],border=GRN,bsz=14,gap=10)
notes(s,"Differentiation and pain peak in the same place. Run gov in parallel for ACV + the SOC2/StateRAMP moat, but software-eng is the velocity/PMF beachhead. ~70/30.")

# 9 — land & expand
s=slide(); kicker(s,"The engine")
title(s,"One data asset, reused three ways.")
card(s,LM,2.35,3.83,4.3,"A · More spend layers",[("Land eng consumption (hardest)",),("→ seats (sales/mktg/support/legal)",),("→ company-wide AI cost system of record",)],border=GRN,bsz=14)
card(s,4.65,2.35,3.83,4.3,"B · More value",[("Attribution → chargeback",),("→ enforcement",),("→ commitment optimization (the prize)",)],border=GRN,bsz=14)
card(s,8.8,2.35,3.93,4.3,"C · More segments",[("Software → enterprise",),("→ financial services",),("→ regulated / gov (metadata-only wedge)",)],border=GRN,bsz=14)
band(s,"Every step reuses the attributed data, adds a buyer or spend layer, raises ACV, deepens lock-in",6.75,0.6,13.5)
notes(s,"Expand the DATA before the features — own the attributed spend so the high-$ commitment optimizer launches from a defensible data advantage.")

# 10 — the prize (now shipped)
s=slide(); pill(s,"The prize · already shipped",w=4.9)
title(s,"Commitment & procurement optimization.",1.35)
bullets(s,LM,2.55,7.0,3.5,[
 "AI is priced like cloud: on-demand vs committed-discount vs provisioned. Most companies over- or under-commit.",
 "Getting it right needs forecasting + attribution — exactly our core data. The highest-$ FinOps lever (cloud analog: 30–60% savings).",
 "Built it: recommend the commit, size it against forfeit risk, pace it, export a vendor negotiation pack — live in-product.",
 "Our work-attributed steadiness signal beats an infra-only optimizer — the moat applied to procurement.",
],sz=15,gap=13)
band(s,"Built before Flexera (ProsperOps + FinOps-for-AI) stitched the model-API layer — from data they don't have",6.4,0.75,14,fill=RGBColor(0xE7,0xF1,0xEC),col=GRND)
notes(s,"The high-$ expansion is no longer a plan — it's shipped (engine + CLI + console: recommender, provisioned break-even, pacing, negotiation pack, advisory caching/batch). Spec: docs/spec-commitment-optimization.md.")

# 11 — TAM
s=slide(); kicker(s,"Market size (US, modeled)")
title(s,"$35–70B governable pool → $1–3B revenue TAM.")
card(s,LM,2.35,3.83,4.3,"Spend under management",[("US governable AI spend pool",),("≈ $35–70B in 2026",),("growing ~50–100%/yr → $100B+",)],border=GRN,bsz=15,gap=11)
card(s,4.65,2.35,3.83,4.3,"Outlay revenue TAM",[("at a 2–5% software take-rate",),("≈ $1.0–3.0B",),("serviceable (ICP) ≈ $0.3–0.9B",)],border=GRN,bsz=15,gap=11)
card(s,8.8,2.35,3.93,4.3,"Core vs incremental",[("Core (attribute+forecast+budget): ~60–70%",),("Incremental (chargeback, enforcement, commitment opt): ~30–40% & rising",)],border=GRN,bsz=14,gap=11)
para(box(s,LM,6.8,CW,0.5),"Demand-side (what companies pay to use AI), not hyperscaler capex. Ranges, not point estimates.",12,MUT,first=True)
notes(s,"Modeled bottoms-up + top-down (market-analysis §5). Capex ($400B+) is the supply-side backdrop, not our pool. Core is the wedge and most of the TAM; incremental features expand the wallet.")

# 12 — why now / plan
s=slide(GRN); kicker(s,"Why now",col=KICKM)
para(box(s,LM,1.6,CW,1.0),"Govern the AI spend before it scales.",36,WHITE,bold=True,font=SERIF,first=True)
para(box(s,LM,2.75,CW,0.8),"The category just inflected, the spend is exploding, and the winning layer (work-attribution) is uncontested. The cheapest time to build the governance layer is now.",18,MINT,first=True)
card(s,LM,3.95,CW,2.75,"The plan",[
 "Now — win the eng-consumption wedge (reposition off 'visibility'); run gov as the compliance-moat lighthouse",
 "Next — expand the account: chargeback + the seat layer (new buyer: Finance)",
 "Done early — the prize is shipped: commitment & procurement optimization, ahead of Flexera; deepen to multi-provider portfolio",
 "Later — broaden segments (enterprise → fin-services → gov-scale) as SOC 2 / StateRAMP land",
],border=KICKM,fill=RGBColor(0x0B,0x51,0x3C),bsz=14,gap=9,tcol=MINT)
notes(s,"Close: the data asset compounds — more spend under management = more lock-in. The high-$ expansion is already built, ahead of schedule. Ask / next step per the audience.")

out=os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),"docs","outlay-strategy-deck.pptx")
prs.save(out); print(f"wrote {out} ({len(prs.slides._sldIdLst)} slides)")
