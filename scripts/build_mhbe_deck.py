#!/usr/bin/env python3
"""Build the MHBE demo deck as a native .pptx (editable in PowerPoint/Keynote/Slides).

All presenter content lives in the deck itself: each slide carries the talking
track + relevant objection answers as PowerPoint speaker notes, and two backup
slides hold the full objection-handling cheat sheet and the security/control detail.
Layouts are sized to fill the 16:9 page edge-to-edge (minimal blank space).

    python scripts/build_mhbe_deck.py   ->   docs/mhbe-demo-deck.pptx
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- brand ---------------------------------------------------------------- #
PAPER = RGBColor(0xFA, 0xF8, 0xF3)
INK   = RGBColor(0x1F, 0x24, 0x30)
MUT   = RGBColor(0x5C, 0x64, 0x70)
GRN   = RGBColor(0x0F, 0x6B, 0x4F)
GRND  = RGBColor(0x0B, 0x51, 0x3C)
GRNL  = RGBColor(0xE7, 0xF1, 0xEC)
AMBER = RGBColor(0x9A, 0x5A, 0x18)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MINT  = RGBColor(0xDD, 0xF0, 0xE7)
KICKM = RGBColor(0xCF, 0xEE, 0xDE)
LINE  = RGBColor(0xE1, 0xDD, 0xD2)
SERIF = "Georgia"
SANS  = "Calibri"
LM    = 0.6           # left margin
CW    = 12.13         # content width (13.333 - 2*0.6)

ASSETS = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs", "assets")
EMUW, EMUH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = EMUW, EMUH
BLANK = prs.slide_layouts[6]


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, EMUW, EMUH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    return tf


def para(tf, text, size, color, *, bold=False, font=SANS, first=False,
         space_after=6, align=PP_ALIGN.LEFT, bullet=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after)
    if bullet:
        text = "•  " + text
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.name = font
    f.color.rgb = color
    return p


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text.strip()


def kicker(s, text, color=GRN, y=0.5):
    para(box(s, LM, y, CW, 0.5), text.upper(), 13, color, bold=True, first=True)


def pill(s, text, x=LM, y=0.5, w=4.7):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.55))
    sh.fill.solid(); sh.fill.fore_color.rgb = GRNL; sh.line.fill.background(); sh.shadow.inherit = False
    sh.text_frame.margin_top = Inches(0.05)
    para(sh.text_frame, text.upper(), 12.5, GRND, bold=True, first=True)


def bullets(s, x, y, w, h, items, size=16, gap=16, color=INK):
    tf = box(s, x, y, w, h)
    for i, t in enumerate(items):
        para(tf, t, size, color, bullet=True, first=(i == 0), space_after=gap)


def card(s, x, y, w, h, title, body_lines, *, border=LINE, fill=WHITE, bsize=14, gap=7):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(1.25); sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.28); tf.margin_right = Inches(0.28)
    tf.margin_top = Inches(0.22); tf.vertical_anchor = MSO_ANCHOR.TOP
    para(tf, title.upper(), 13, MUT, bold=True, first=True, space_after=10)
    for ln in body_lines:
        para(tf, ln, bsize, INK, bullet=True, space_after=gap)
    return sh


def pic(s, name, x, y, w, border=True, bcolor=None):
    path = os.path.join(ASSETS, name)
    iw, ih = Image.open(path).size
    p = s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    if border:
        p.line.color.rgb = bcolor or LINE; p.line.width = Pt(1.25)
    return w * ih / iw


# ========================================================================== #
# 1 — title
s = slide(); kicker(s, "Outlay · for Maryland Health Connection")
tf = box(s, LM, 1.45, CW, 3)
para(tf, "AI spend, mapped to the work", 46, INK, bold=True, font=SERIF, first=True, space_after=2)
para(tf, "that drove it — and held to budget.", 46, GRN, bold=True, font=SERIF, space_after=16)
para(tf, "A read-only, metadata-only governance layer for the AI behind Maryland Health Connection.", 22, MUT)
sh = s.shapes.add_shape(1, Inches(LM), Inches(5.55), Inches(CW), Inches(0.9))
sh.fill.solid(); sh.fill.fore_color.rgb = GRNL; sh.line.fill.background(); sh.shadow.inherit = False
sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
para(sh.text_frame, "Attribute  ·  Forecast  ·  Govern", 24, GRND, bold=True, font=SERIF, first=True, align=PP_ALIGN.CENTER)
para(box(s, LM, 6.85, CW, 0.4), "Outlay.ai   ·   Demo · June 26, 2026 · [presenter]", 12, MUT, bold=True, first=True)
notes(s, """OPEN HERE. The one rule for this room: lead with the architecture, not the features.
This is a state health exchange — their first instinct is "what data does this touch?" Win that in
the first five minutes (slide 4) and the rest lands. Room: CIO (Koshanam)=technical fit,
Compliance (Brennan)=the gate, CFO (Armiger)=budget value. Speak to all three.""")

# 2 — the moment
s = slide(); kicker(s, "The moment")
para(box(s, LM, 1.35, CW, 1.6), "Maryland is operationalizing its 2025 AI Enablement Strategy.",
     38, INK, bold=True, font=SERIF, first=True)
para(box(s, LM, 3.25, CW, 0.6), "Right after “is it safe?” comes the question every CFO and CIO asks:",
     22, MUT, first=True)
sh = s.shapes.add_shape(1, Inches(LM), Inches(4.05), Inches(CW), Inches(1.25))
sh.fill.solid(); sh.fill.fore_color.rgb = GRNL; sh.line.fill.background(); sh.shadow.inherit = False
sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
para(sh.text_frame, "“What is the AI costing us — per program — and is it on budget?”",
     30, GRND, bold=True, font=SERIF, first=True, align=PP_ALIGN.CENTER)
para(box(s, LM, 5.7, CW, 1.2), "Today that answer lives in invoices and spreadsheets, after the money is "
     "spent. Outlay answers it continuously, attributed to the work.", 19, MUT, first=True)
notes(s, """Talking track: "Maryland's standing up its AI strategy this winter. The cost-and-governance
question always lands right after the safety question — and right now it's answered in arrears, in
invoices. We make it continuous and attributed." Keep it short; this is the setup for slide 4.""")

# 3 — what it does
s = slide(); kicker(s, "What Outlay does")
para(box(s, LM, 1.05, CW, 0.8), "Three things, from data you already have.", 32, INK, bold=True, font=SERIF, first=True)
bullets(s, LM, 2.35, 4.7, 4.4, [
    "Attribute — every AI dollar mapped to a team, work type, ticket, and engineer.",
    "Forecast — what planned work will cost, back-tested on your own delivered work.",
    "Govern — budgets per program, real-time pacing, and on-/off-track alerts.",
], size=19, gap=22)
para(box(s, LM, 6.5, 4.7, 0.6), "Read-only links — no agents, nothing in the request path.", 13, MUT, first=True)
pic(s, "shot-estimate.png", 5.55, 1.7, 7.2)
notes(s, """Talking track: "From two read-only connections — your tracker and your AI provider's billing
API — Outlay does attribution, forecasting, and governance. No agents, nothing inline." Don't dwell;
the architecture slide next is where you slow down.""")

# 4 — architecture (the compliance slide)
s = slide(); pill(s, "Why this is safe for a health exchange")
para(box(s, LM, 1.3, CW, 0.9), "We never see PHI, PII, or FTI — by design.", 32, INK, bold=True, font=SERIF, first=True)
bullets(s, LM, 2.45, 4.35, 4.4, [
    "Metadata-only — token counts, ticket IDs, work types, dollar figures. Never prompts or outputs.",
    "Bring-your-own-key — your API keys stay in your environment.",
    "Read-only — never a proxy or gateway in any data path.",
    "Out of MARS-E / IRS 1075 / HIPAA data scope — what stops most vendors at the door.",
], size=16, gap=18)
h = pic(s, "diagram-architecture.png", 5.15, 2.35, 7.85)
notes(s, """THE slide. Slow down. "Everything I'll show runs on metadata only — token counts, ticket
IDs, dollars. Never prompts, outputs, PHI, PII, or FTI. Your keys stay with you; we're read-only,
never in a data path. So none of the data MARS-E, IRS 1075, or HIPAA govern reaches us — the
integration is outside that boundary."
OBJECTION — "Do you ever see PHI/member data?": "No, by design. Metadata only. Prompts and keys
never leave your environment."  Name MARS-E / ARC-AMPE / 1075 yourself — it signals you know their world.""")

# 5 — demo divider
s = slide(GRN); kicker(s, "Live", color=KICKM)
para(box(s, LM, 1.9, 5.3, 1.5), "Let's look at the product.", 38, WHITE, bold=True, font=SERIF, first=True)
para(box(s, LM, 3.5, 5.2, 2.5), "On sample data shaped like a real engineering program — the same screens "
     "you'd see on day one of a pilot.", 19, MINT, first=True)
pic(s, "shot-overview.png", 6.25, 1.5, 6.75, bcolor=WHITE)
notes(s, """SETUP (do before the call): signed into a demo-flagged account on app.outlay-ai.com with
SAMPLE DATA loaded; "business" persona (or "eng" if technical); ~110% zoom; notifications off.
Smoke it first: python scripts/preflight.py --base https://app.outlay-ai.com, then click every page once.
Backup: localhost — CONSOLE_SECRET=… DEMO_ACCOUNT_EMAILS="*" python -m console.server.""")


def demo_slide(kick, title, blist, shot, notes_text):
    s = slide(); kicker(s, kick)
    para(box(s, LM, 1.05, 4.05, 1.5), title, 26, INK, bold=True, font=SERIF, first=True)
    bullets(s, LM, 2.7, 4.0, 3.8, blist, size=15, gap=16)
    pic(s, shot, 4.75, 1.05, 8.3)
    notes(s, notes_text)
    return s

# 6 — attribution
demo_slide("Demo · Attribution", "Every dollar, on the roadmap.", [
    "By team & cost center, work type, ticket, and engineer.",
    "Ticket coverage shown honestly — and how to lift it.",
    "FOCUS-aligned CSV export for finance.",
], "shot-spend.png", """Click path: open Spend (/app/outlay). Scroll the attribution (team / work type / ticket /
engineer). Then point at the COVERAGE line: "We're honest about it — this is the share of spend that
maps to a specific work item, and exactly how to lift it. We never inflate the number." Mention the
FOCUS CSV export ("the artifact your finance team loads").""")

# 7 — accuracy
demo_slide("Demo · Accuracy (the honesty layer)", "We don't ask you to trust a vendor benchmark.", [
    "Back-tested on your own closed tickets (leave-one-out).",
    "Median error, within-P90, and sample size — never hidden.",
    "Per-work-type accuracy with over/under bias.",
], "shot-accuracy.png", """THE TRUST MOMENT — this is what separates you from a dashboard. Open Accuracy
(/app/outlay/accuracy). "The #1 question is 'can I trust the forecast?' We don't answer with a vendor
benchmark — we back-test on YOUR closed tickets, leave-one-out, and show the measured error, sample
size never hidden." Let a compliance/finance person sit with this; candor is the sell.""")

# 8 — governance
demo_slide("Demo · Governance", "Put AI compute on a budget — per program.", [
    "Program budgets across teams, projects, and work types.",
    "Real-time pacing → the projected breach date.",
    "On-/off-track earned-value rating; alerts to your SIEM.",
], "shot-governance.png", """The CFO segment. Open Budgets/Governance (/app/outlay/budgets, /app/outlay/governance).
"Set a budget per program. We pace it in real time, project the breach date, and rate it on-track /
watch / off-track from forecast vs actual on completed work — and it reaches you proactively:
digest, close pack, Slack, webhook to your SIEM."  """)

# 9 — security posture
s = slide(); kicker(s, "Security & compliance posture")
para(box(s, LM, 1.1, CW, 0.8), "Straight talk on where we are.", 32, INK, bold=True, font=SERIF, first=True)
card(s, LM, 2.3, 5.9, 4.45, "In place today",
     ["SSO (OIDC) + SCIM provisioning", "Phishing-resistant MFA (passkeys)",
      "RBAC · audit log → SIEM export", "Retention + self-serve erasure",
      "Encryption at rest + KMS hook", "IR plan · WCAG 2.1 AA · CI security scan"],
     border=GRN, bsize=15, gap=11)
card(s, 6.83, 2.3, 5.9, 4.45, "On the roadmap (honest)",
     ["SOC 2 Type II — in progress", "Annual penetration test — scheduled",
      "StateRAMP / GovRAMP — roadmap", "FedRAMP cloud + FIPS crypto — deal-gated",
      "Metadata-only keeps most of these out of scope for the data we touch."],
     border=AMBER, bsize=15, gap=11)
notes(s, """OBJECTION HANDLING — keep these cold:
• "SOC 2 / StateRAMP / FedRAMP?" → "SOC 2 Type II in progress, StateRAMP on the roadmap. I'll be
  straight — not there yet. What makes a pilot work NOW is the architecture: read-only + metadata-only
  puts us outside the data scope those frameworks govern. Happy to share our SOC 2 timeline + NIST
  800-53 mapping today."
• "MARS-E / ARC-AMPE / 1075?" → "Those govern systems that receive member data/FTI. We receive none,
  so we're outside that boundary. For controls that apply to any vendor — access, audit, encryption,
  IR — we have them."
• "Where are you hosted?" → "Today Fly.io. For a gov engagement that advances, the path is a re-host
  to AWS GovCloud / Azure Government — which also gives FIPS-validated crypto + US residency. Gated on
  a real deal; metadata-only lets us pilot before it."
NEVER say "HIPAA compliant" or claim a cert you don't hold. Candor wins with a compliance officer.""")

# 10 — fit
s = slide(); kicker(s, "Fit for MHBE")
para(box(s, LM, 1.1, CW, 0.8), "Two ways this earns its place now.", 32, INK, bold=True, font=SERIF, first=True)
card(s, LM, 2.3, 5.9, 4.45, "The AI-governance layer",
     ["Be the cost-attribution + budget-guardrail layer as you stand up the 2025 AI strategy.",
      "Governance-first — matched to how your program is being built.",
      "Day-one value with no instrumentation: read-only connections."],
     border=GRN, bsize=16, gap=14)
card(s, 6.83, 2.3, 5.9, 4.45, "Vendor & enhancement spend",
     ["As Deloitte / J29 enhancement work adopts AI coding agents, see cost per project.",
      "The earned-value view a public agency needs to justify IT spend.",
      "Attributes the $12.7M “MHC enhancements” work to the tickets behind it."],
     border=GRN, bsize=16, gap=14)
notes(s, """Tie it to THEIR world. The $12.7M "MHC enhancements" line + the Deloitte/J29 contracts are
the natural first surface — AI-assisted dev cost per project, which a public agency must justify.""")

# 11 — the ask
s = slide(); pill(s, "The ask", w=2.3)
para(box(s, LM, 1.35, CW, 0.9), "A two-week, read-only pilot — zero data-handling risk.",
     32, INK, bold=True, font=SERIF, first=True)
bullets(s, LM, 2.75, CW, 3.4, [
    "Scoped to your development / vendor AI usage, not any production PHI system.",
    "Metadata-only and read-only — nothing for the security review to clear on data flow.",
    "You walk away with a real number: attribution coverage + a measured forecast accuracy on your work.",
], size=20, gap=22)
sh = s.shapes.add_shape(1, Inches(LM), Inches(5.95), Inches(CW), Inches(0.85))
sh.fill.solid(); sh.fill.fore_color.rgb = GRNL; sh.line.fill.background(); sh.shadow.inherit = False
sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
para(sh.text_frame, "Free during the pilot · we share our SOC 2 timeline and NIST 800-53 control mapping up front.",
     16, GRND, bold=True, first=True, align=PP_ALIGN.CENTER)
notes(s, """Don't leave without a next step. Say it plainly (above). Then pin: a follow-up with Scott
Brennan on the architecture + control mapping, and identify the dev/vendor data sources for the pilot.
OBJECTION — "What does a pilot require from us?" → "A read-only token to a work tracker and your AI
usage/billing API, scoped to dev/vendor. No installs, no production access, nothing inline."
OBJECTION — "Cost?" → "Pilot is free. Platform pricing we set with early customers — but prove the
number first." Leave-behind = the one-pager on the next slide's content.""")

# 12 — close
s = slide(GRN); kicker(s, "Why now", color=KICKM)
para(box(s, LM, 2.05, CW, 1.6), "Govern the AI spend before it scales.", 46, WHITE, bold=True, font=SERIF, first=True)
para(box(s, LM, 3.95, CW, 1.6), "The cheapest time to put AI on a budget is at the start of the curve — "
     "which is exactly where Maryland is.", 22, MINT, first=True)
para(box(s, LM, 6.85, CW, 0.4), "[presenter] · [email] · outlay-ai.com", 13, RGBColor(0xBF,0xE0,0xD0), first=True)
notes(s, "Close, then go to the ask (slide 11) and agree the next step. Backup slides follow for Q&A.")

# 13 — BACKUP: objection handling
s = slide(); kicker(s, "Backup · Objection handling", color=MUT)
tf = box(s, LM, 1.15, CW, 5.9)
qa = [
 ("Do you ever see PHI / member data / prompts?", "No — by design. Metadata only: token counts, ticket IDs, work types, dollars. Prompts, outputs, and API keys never leave your environment."),
 ("SOC 2 / StateRAMP / FedRAMP authorized?", "SOC 2 Type II in progress; StateRAMP on the roadmap. Not there yet — and the architecture is what makes a pilot work now (outside the data scope those govern)."),
 ("MARS-E / ARC-AMPE / IRS 1075?", "Those govern systems that receive member data/FTI. We receive none, so we're outside that boundary. Vendor-applicable controls (access/audit/encryption/IR) we have."),
 ("Where are you hosted?", "Fly.io today. A gov engagement that advances → re-host to AWS GovCloud / Azure Gov (FIPS crypto + US residency). Gated on a real deal; metadata-only lets us pilot first."),
 ("What does a pilot require from us?", "A read-only token to a tracker + your AI usage/billing API, scoped to dev/vendor. No installs, no production access, nothing inline. Two weeks."),
 ("Cost?", "Pilot is free. Platform pricing set with early customers — prove the number first."),
]
for i, (q, a) in enumerate(qa):
    para(tf, q, 16, GRN, bold=True, first=(i == 0), space_after=3)
    para(tf, a, 14, INK, space_after=13)
notes(s, "Flip here if Q&A goes deep on compliance. Never overclaim — 'in progress / roadmap / out of scope by design'.")

# 14 — BACKUP: control summary
s = slide(); kicker(s, "Backup · Security & control summary", color=MUT)
para(box(s, LM, 1.05, CW, 0.9), "Architecture removes the heaviest scope; operational controls are in place; "
     "attestations are the roadmap.", 19, INK, bold=True, font=SERIF, first=True)
card(s, LM, 2.15, 5.9, 4.7, "Have & can evidence",
     ["Metadata-only / BYOK / read-only — no PHI/PII/FTI", "SSO (OIDC) + SCIM",
      "MFA incl. phishing-resistant passkeys (AAL2/3)", "RBAC · audit log + SIEM export (/api/v1/audit)",
      "Session controls: idle/absolute/epoch revocation", "Retention + self-serve erasure",
      "Encryption at rest + pluggable KMS hook", "IR plan (MD-SOC 1-hr) + incident webhook",
      "WCAG 2.1 AA + automated a11y & security CI scans"], border=GRN, bsize=13.5, gap=7)
card(s, 6.83, 2.15, 5.9, 4.7, "Roadmap / external / funded",
     ["SOC 2 Type II — in progress (highest leverage)", "Annual third-party penetration test",
      "Monthly vulnerability-scan program + patch SLA", "StateRAMP / GovRAMP (Ready → Authorized)",
      "FedRAMP-authorized cloud (AWS GovCloud / Azure Gov)", "FIPS 140-validated crypto (IRS 1075)",
      "3PAO assessment vs NIST 800-53 Rev 5"], border=AMBER, bsize=13.5, gap=7)
notes(s, "Mirrors docs/prospect-maryland-health-connection.md (the requirements→readiness matrix) and docs/soc2-stateramp-sequencing.md.")

out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs", "mhbe-demo-deck.pptx")
prs.save(out)
print(f"wrote {out}  ({len(prs.slides._sldIdLst)} slides)")
